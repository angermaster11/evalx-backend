from langchain.chat_models import ChatOpenAI
from langchain.schema import HumanMessage, SystemMessage
from pptx import Presentation
import base64, io, requests, json, re, httpx
from dotenv import load_dotenv
from api.models.state import State
from fastapi import APIRouter

load_dotenv()


# Initialize GPT-4o with deterministic output
llm = ChatOpenAI(model="gpt-4o", temperature=0)


async def extract_ppt_text(file_path: str) -> str:
    """Extracts plain text from a PPT (ignores per-slide details)."""
    if file_path.startswith("http"):
        async with httpx.AsyncClient() as client:
            response = await client.get(file_path)
        prs = Presentation(io.BytesIO(response.content))
    else:
        prs = Presentation(file_path)

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())

    return "\n".join(all_text)


def clean_json_output(raw_text: str):
    """Cleans model output to extract valid JSON."""
    cleaned = re.sub(r"^```json|```$", "", raw_text.strip(), flags=re.MULTILINE).strip()
    try:
        return json.loads(cleaned)
    except Exception:
        return {"error": "Model did not return valid JSON", "raw": raw_text}



async def analyze_ppt_with_gpt(state: State) -> State:
    ppt_text = await extract_ppt_text(state["file_path"])

    messages = [
        SystemMessage(content=(
            "You are a strict PPT Evaluator.\n"
            "Analyze the slides and respond ONLY in JSON.\n"
            "Use the provided hackathon description to contextualize your feedback.\n"
            "The JSON schema is:\n"
            "{\n"
            '  "impact": float (0-1),\n'
            '  "uniqueness": float (0-1),\n'
            '  "clarity": float (0-1),\n'
            '  "design": float (0-1),\n'
            '  "advantages": [string],\n'
            '  "disadvantages": [string],\n'
            '  "improvements": [string],\n'
            '  "overall_report": string,\n'
            '  "overall_feedback": string,\n'
            '  "loose_points": string\n'
            "}\n"
            "⚠️ Important: Do not include markdown fences, only raw JSON."
        )),
        HumanMessage(content=(
            f"Hackathon Context:\n{state['content']}\n\n"
            f"PPT Content:\n{ppt_text}"
        ))
    ]

    # ✅ Await the LLM call
    response = await llm.ainvoke(messages)

    result = clean_json_output(response.content)
    return {
        **state,
        "output": result
    }


